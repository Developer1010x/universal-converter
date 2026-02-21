#!/usr/bin/env python3
"""
Universal Converter - Convert Anything to Anything
OS-neutral, minimal dependencies, comprehensive conversion tool
"""

import os
import sys
import json
import csv
import xml.etree.ElementTree as ET
import base64
import hashlib
import shutil
import tempfile
import subprocess
import logging
import struct
import zipfile
import tarfile
import configparser
import re
import mimetypes
from pathlib import Path
from typing import Union, Dict, Any, Optional, List, Tuple
from datetime import datetime
from dataclasses import dataclass, field
import traceback
import io
import stat


logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


REQUIRED_PACKAGES = {
    'pillow': 'Pillow',
    'pypdf2': 'PyPDF2',
    'reportlab': 'reportlab',
    'python-docx': 'docx',
    'pandas': 'pandas',
    'openpyxl': 'openpyxl',
    'xlrd': 'xlrd',
    'python-pptx': 'pptx',
    'markdown': 'markdown',
    'weasyprint': 'weasyprint',
    'odfpy': 'odf',
    'cairosvg': 'cairosvg',
    'pydub': 'pydub',
    'speechrecognition': 'SpeechRecognition',
    'moviepy': 'moviepy',
    'rarfile': 'rarfile',
    'pyzipper': 'pyzipper',
    'beautifulsoup4': 'bs4',
    'html2text': 'html2text',
    'requests': 'requests',
    'xmltodict': 'xmltodict',
    'numpy': 'numpy',
    'scipy': 'scipy',
    'h5py': 'h5py',
    'pyarrow': 'pyarrow',
    'qrcode': 'qrcode',
    'pyzbar': 'pyzbar',
    'bcrypt': 'bcrypt',
    'argon2': 'argon2',
    'pyjwt': 'jwt',
    'ebooklib': 'ebooklib',
    'pint': 'pint',
    'forex-python': 'forex_python',
    'pytesseract': 'pytesseract',
    'pdf2image': 'pdf2image',
    'msgpack': 'msgpack',
}


def ensure_dependencies():
    """Check and install missing dependencies"""
    import importlib
    
    missing = []
    for pkg, import_name in REQUIRED_PACKAGES.items():
        try:
            importlib.import_module(import_name)
        except ImportError:
            missing.append(pkg)
    
    if missing:
        logger.info(f"Installing missing dependencies: {', '.join(missing)}")
        try:
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', *missing])
        except subprocess.CalledProcessError:
            logger.warning(f"Failed to install some packages. Install manually: pip install {' '.join(missing)}")


def ensure_package(package_name: str):
    """Ensure a specific package is installed"""
    import importlib
    
    pkg_info = REQUIRED_PACKAGES.get(package_name.lower())
    if not pkg_info:
        pkg_info = package_name
    
    try:
        importlib.import_module(pkg_info)
        return True
    except ImportError:
        try:
            logger.info(f"Installing {package_name}...")
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', package_name])
            return True
        except subprocess.CalledProcessError:
            return False


# Auto-check dependencies on import
try:
    ensure_dependencies()
except Exception:
    pass


class ConversionError(Exception):
    def __init__(self, message: str, source_format: str = None, target_format: str = None, details: str = None):
        self.source_format = source_format
        self.target_format = target_format
        self.details = details
        super().__init__(message)
    
    def __str__(self):
        msg = super().__str__()
        if self.source_format and self.target_format:
            msg = f"{msg} ({self.source_format} → {self.target_format})"
        return msg


@dataclass
class ConversionResult:
    success: bool
    data: Any = None
    output_path: str = None
    error: str = None
    metadata: Dict = field(default_factory=dict)
    warnings: List[str] = field(default_factory=list)


class UniversalConverter:
    """Main universal converter - OS neutral, minimal dependencies"""
    
    TEXT_FORMATS = {'json', 'csv', 'xml', 'yaml', 'yml', 'txt', 'html', 'htm', 'markdown', 'md', 'rst', 'sql', 'ini', 'toml', 'url', 'log'}
    DATA_FORMATS = {'json', 'csv', 'tsv', 'xml', 'yaml', 'yml', 'ini', 'toml', 'cfg', 'conf'}
    DOCUMENT_FORMATS = {'pdf', 'docx', 'doc', 'odt', 'rtf', 'tex', 'epub', 'mobi'}
    IMAGE_FORMATS = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif', 'webp', 'ico', 'ppm', 'pgm', 'pbm', 'jp2', 'svg'}
    ARCHIVE_FORMATS = {'zip', 'tar', 'gz', 'bz2', 'xz', '7z', 'rar'}
    CODE_FORMATS = {'py', 'js', 'ts', 'java', 'c', 'cpp', 'h', 'hpp', 'cs', 'go', 'rs', 'rb', 'php', 'swift', 'kt', 'scala', 'r', 'sql', 'sh', 'bash', 'ps1'}
    
    DATABASE_FORMATS = {'mysql', 'postgresql', 'postgres', 'sqlite', 'mongodb', 'mongo', 'redis', 'sql', 'json', 'csv', 'parquet', 'db'}
    
    ALL_FORMATS = TEXT_FORMATS | DATA_FORMATS | DOCUMENT_FORMATS | IMAGE_FORMATS | ARCHIVE_FORMATS | CODE_FORMATS | DATABASE_FORMATS
    
    def __init__(self):
        self.errors: List[str] = []
        self.warnings: List[str] = []
    
    def detect_format(self, data: Any) -> Optional[str]:
        if isinstance(data, str):
            path = Path(data)
            if path.exists() and path.is_file():
                ext = path.suffix[1:].lower()
                if ext:
                    return ext
                try:
                    with open(path, 'r', encoding='utf-8') as f:
                        content = f.read(1000)
                    return self._detect_text_format(content)
                except:
                    return None
        if isinstance(data, dict): return 'json'
        if isinstance(data, list): return 'json'
        return None
    
    def _detect_text_format(self, content: str) -> str:
        content = content.strip()
        if content.startswith('{') or content.startswith('['): return 'json'
        if content.startswith('<?xml'): return 'xml'
        if '<!DOCTYPE html>' in content or '<html' in content.lower(): return 'html'
        if content.startswith('#'): return 'markdown'
        if '---' in content and ':' in content: return 'yaml'
        if content.startswith('<?'): return 'xml'
        if ',' in content and '\n' in content: return 'csv'
        return 'txt'
    
    def convert(self, data: Any, from_fmt: str = None, to_fmt: str = None, output_path: str = None, **options) -> ConversionResult:
        if from_fmt is None:
            from_fmt = self.detect_format(data)
        
        if from_fmt is None or to_fmt is None:
            return ConversionResult(success=False, error="Could not detect format")
        
        from_fmt, to_fmt = from_fmt.lower(), to_fmt.lower()
        
        try:
            if isinstance(data, str) and Path(data).exists() and Path(data).is_file():
                return self.convert_file(data, output_path or f"output.{to_fmt}", from_fmt, to_fmt)
            
            result = self._convert_data(data, from_fmt, to_fmt, options)
            
            if output_path:
                Path(output_path).parent.mkdir(parents=True, exist_ok=True)
                mode = 'wb' if isinstance(result, bytes) else 'w'
                encoding = None if isinstance(result, bytes) else 'utf-8'
                with open(output_path, mode, encoding=encoding) as f:
                    f.write(result if isinstance(result, (bytes, str)) else str(result))
                return ConversionResult(success=True, output_path=output_path, metadata={'format': to_fmt})
            
            return ConversionResult(success=True, data=result, metadata={'format': to_fmt})
            
        except Exception as e:
            return ConversionResult(success=False, error=str(e), details=traceback.format_exc())
    
    def _convert_data(self, data: Any, from_fmt: str, to_fmt: str, options: Dict) -> Any:
        data = self._parse_input(data, from_fmt)
        
        converters = {
            'json': self._to_json,
            'csv': self._to_csv,
            'tsv': self._to_tsv,
            'xml': self._to_xml,
            'yaml': self._to_yaml,
            'yml': self._to_yaml,
            'txt': self._to_text,
            'html': self._to_html,
            'markdown': self._to_markdown,
            'md': self._to_markdown,
            'sql': self._to_sql,
            'ini': self._to_ini,
            'toml': self._to_toml,
            'url': self._to_urlencoded,
            'base64': self._to_base64,
            'query': self._to_sql_query,
        }
        
        if to_fmt in converters:
            return converters[to_fmt](data, options)
        
        raise ConversionError(f"Unsupported: {from_fmt} → {to_fmt}", from_fmt, to_fmt)
    
    def _parse_input(self, data: Any, fmt: str) -> Dict:
        if isinstance(data, dict): return data
        if isinstance(data, list): return {'data': data}
        if isinstance(data, str):
            try:
                return json.loads(data)
            except:
                return {'text': data}
        return {'value': str(data)}
    
    def _to_json(self, data: Dict, opts: Dict) -> str:
        return json.dumps(data, indent=opts.get('indent', 2), ensure_ascii=False, default=str)
    
    def _to_csv(self, data: Dict, opts: Dict) -> str:
        if isinstance(data, dict):
            if all(isinstance(v, list) for v in data.values()):
                return self._dict_of_lists_to_csv(data)
            data = [data]
        if not isinstance(data, list): data = [data]
        if not data: return ""
        
        keys = set()
        for item in data:
            if isinstance(item, dict): keys.update(item.keys())
        if not keys: return ""
        
        keys = list(keys)
        rows = [','.join(keys)]
        for item in data:
            if isinstance(item, dict):
                rows.append(','.join(self._escape_csv(str(item.get(k, ''))) for k in keys))
        return '\n'.join(rows)
    
    def _escape_csv(self, val: str) -> str:
        if ',' in val or '"' in val or '\n' in val:
            return f'"{val.replace("\"", "\"\"")}"'
        return val
    
    def _dict_of_lists_to_csv(self, data: Dict) -> str:
        max_len = max(len(v) for v in data.values()) if data else 0
        keys = list(data.keys())
        rows = [','.join(keys)]
        for i in range(max_len):
            rows.append(','.join(str(data[k][i] if i < len(data[k]) else '') for k in keys))
        return '\n'.join(rows)
    
    def _to_tsv(self, data: Dict, opts: Dict) -> str:
        csv_str = self._to_csv(data, opts)
        return csv_str.replace(',', '\t')
    
    def _to_xml(self, data: Any, opts: Dict) -> str:
        root = ET.Element('root')
        self._dict_to_xml_element(data, root)
        return '<?xml version="1.0" encoding="UTF-8"?>\n' + ET.tostring(root, encoding='unicode')
    
    def _dict_to_xml_element(self, data: Any, parent: ET.Element):
        if isinstance(data, dict):
            for k, v in data.items():
                child = ET.SubElement(parent, str(k))
                self._dict_to_xml_element(v, child)
        elif isinstance(data, list):
            for item in data:
                child = ET.SubElement(parent, 'item')
                self._dict_to_xml_element(item, child)
        else:
            parent.text = str(data) if data is not None else ''
    
    def _to_yaml(self, data: Any, opts: Dict) -> str:
        return self._yaml_convert(data, 0)
    
    def _yaml_convert(self, data: Any, indent: int) -> str:
        lines = []
        prefix = '  ' * indent
        
        if isinstance(data, dict):
            for k, v in data.items():
                if isinstance(v, (dict, list)):
                    lines.append(f"{prefix}{k}:")
                    lines.append(self._yaml_convert(v, indent + 1))
                else:
                    lines.append(f"{prefix}{k}: {self._yaml_escape(v)}")
        elif isinstance(data, list):
            for item in data:
                if isinstance(item, (dict, list)):
                    lines.append(f"{prefix}-")
                    lines.append(self._yaml_convert(item, indent + 1))
                else:
                    lines.append(f"{prefix}- {self._yaml_escape(item)}")
        else:
            lines.append(f"{prefix}{self._yaml_escape(data)}")
        
        return '\n'.join(lines)
    
    def _yaml_escape(self, v: Any) -> str:
        if v is None: return 'null'
        if isinstance(v, bool): return 'true' if v else 'false'
        v = str(v)
        if any(c in v for c in ':#[]{},&*!|>\'"%@`') or v.startswith('- ') or v.startswith(' '):
            return f'"{v.replace("\\", "\\\\").replace("\"", "\\\"")}"'
        return v
    
    def _to_text(self, data: Any, opts: Dict) -> str:
        lines = []
        if isinstance(data, dict):
            for k, v in data.items():
                lines.append(f"{k}: {v}")
        elif isinstance(data, list):
            for item in data:
                lines.append(str(item))
        else:
            return str(data)
        return '\n'.join(lines)
    
    def _to_html(self, data: Any, opts: Dict) -> str:
        html = ['<!DOCTYPE html>', '<html>', '<head>',
                '<meta charset="UTF-8"><title>Converted</title>',
                '<style>',
                'body{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,sans-serif;max-width:900px;margin:0 auto;padding:20px;line-height:1.6}',
                'table{border-collapse:collapse;width:100%;margin:20px 0}',
                'th,td{border:1px solid #ddd;padding:12px;text-align:left}',
                'th{background:#f5f5f5}',
                'code{background:#f4f4f4;padding:2px 6px;border-radius:3px}',
                'pre{background:#f4f4f5;padding:15px;border-radius:6px;overflow-x:auto}',
                '</style></head><body>']
        
        if isinstance(data, dict):
            html.append('<table>')
            for k, v in data.items():
                html.append(f'<tr><th>{k}</th><td>{self._escape_html(v)}</td></tr>')
            html.append('</table>')
        elif isinstance(data, list):
            html.append('<ul>')
            for item in data:
                html.append(f'<li>{self._escape_html(item)}</li>')
            html.append('</ul>')
        else:
            html.append(f'<p>{self._escape_html(data)}</p>')
        
        html.extend(['</body>', '</html>'])
        return '\n'.join(html)
    
    def _escape_html(self, s: Any) -> str:
        s = str(s)
        return s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
    
    def _to_markdown(self, data: Any, opts: Dict) -> str:
        lines = ['# Converted Document', '']
        
        if isinstance(data, dict):
            for k, v in data.items():
                lines.append(f"## {k}")
                lines.append(self._md_value(v))
                lines.append('')
        elif isinstance(data, list):
            for item in data:
                if isinstance(item, dict):
                    for k, v in item.items():
                        lines.append(f"- **{k}**: {v}")
                else:
                    lines.append(f"- {item}")
            lines.append('')
        else:
            lines.append(str(data))
        
        return '\n'.join(lines)
    
    def _md_value(self, data: Any, indent: int = 0) -> str:
        lines = []
        prefix = '  ' * indent
        
        if isinstance(data, dict):
            for k, v in data.items():
                if isinstance(v, (dict, list)):
                    lines.append(f"{prefix}- **{k}**:")
                    lines.append(self._md_value(v, indent + 1))
                else:
                    lines.append(f"{prefix}- **{k}**: {v}")
        elif isinstance(data, list):
            for item in data:
                lines.append(f"{prefix}- {item}")
        
        return '\n'.join(lines)
    
    def _to_sql(self, data: Any, opts: Dict) -> str:
        lines = []
        table = opts.get('table', 'data')
        
        if isinstance(data, dict) and 'tables' in data:
            for tbl, tbl_data in data['tables'].items():
                cols = tbl_data.get('columns', [])
                rows = tbl_data.get('data', [])
                col_defs = ', '.join([f"{c} TEXT" for c in cols]) if cols else "id INTEGER PRIMARY KEY"
                lines.append(f"CREATE TABLE {tbl} ({col_defs});")
                for row in rows:
                    if isinstance(row, dict):
                        vals = "', '".join(str(v) for v in row.values())
                        lines.append(f"INSERT INTO {tbl} VALUES ('{vals}');")
        else:
            if not isinstance(data, list): data = [data]
            if data and isinstance(data[0], dict):
                cols = list(data[0].keys())
                col_defs = ', '.join([f"{c} TEXT" for c in cols])
                lines.append(f"CREATE TABLE {table} ({col_defs});")
                for row in data:
                    if isinstance(row, dict):
                        vals = "', '".join(str(v) for v in row.values())
                        lines.append(f"INSERT INTO {table} VALUES ('{vals}');")
        
        return '\n'.join(lines)
    
    def _to_ini(self, data: Dict, opts: Dict) -> str:
        config = configparser.ConfigParser()
        
        if isinstance(data, dict):
            if all(isinstance(v, dict) for v in data.values()):
                for section, values in data.items():
                    config[section] = values
            else:
                config['DEFAULT'] = data
        
        out = io.StringIO()
        config.write(out)
        return out.getvalue()
    
    def _to_toml(self, data: Any, opts: Dict) -> str:
        return self._to_yaml(data, opts)
    
    def _to_urlencoded(self, data: Any, opts: Dict) -> str:
        from urllib.parse import urlencode
        flat = self._flatten_dict(data) if isinstance(data, dict) else {'data': str(data)}
        return urlencode(flat)
    
    def _flatten_dict(self, d: Dict, parent_key: str = '', sep: str = '_') -> Dict:
        items = []
        for k, v in d.items():
            new_key = f"{parent_key}{sep}{k}" if parent_key else k
            if isinstance(v, dict):
                items.extend(self._flatten_dict(v, new_key, sep=sep).items())
            else:
                items.append((new_key, str(v)))
        return dict(items)
    
    def _to_base64(self, data: Any, opts: Dict) -> str:
        return base64.b64encode(str(data).encode('utf-8')).decode('utf-8')
    
    def _to_sql_query(self, data: Any, opts: Dict) -> str:
        queries = []
        if not isinstance(data, list): data = [data]
        for item in data:
            if isinstance(item, dict):
                conds = ' AND '.join(f"{k} = '{v}'" for k, v in item.items())
                queries.append(f"SELECT * FROM table WHERE {conds};")
        return '\n'.join(queries)
    
    def convert_file(self, input_path: str, output_path: str = None, from_fmt: str = None, to_fmt: str = None) -> ConversionResult:
        path = Path(input_path)
        
        if not path.exists():
            return ConversionResult(success=False, error=f"File not found: {input_path}")
        
        if from_fmt is None:
            from_fmt = path.suffix[1:].lower()
        
        if output_path is None:
            output_path = str(path.parent / f"{path.stem}.converted")
        
        if to_fmt is None:
            to_fmt = Path(output_path).suffix[1:].lower()
        
        try:
            if from_fmt == to_fmt:
                shutil.copy2(input_path, output_path)
                return ConversionResult(success=True, output_path=output_path)
            
            if from_fmt in self.IMAGE_FORMATS and to_fmt in self.IMAGE_FORMATS:
                return self._convert_image(input_path, output_path, from_fmt, to_fmt)
            
            if from_fmt == 'pdf' and to_fmt == 'txt':
                return self._pdf_to_text(input_path, output_path)
            
            if from_fmt == 'txt' and to_fmt == 'pdf':
                return self._text_to_pdf(input_path, output_path)
            
            if from_fmt == 'pdf' and to_fmt in ('docx', 'doc'):
                return self._convert_pdf_word(input_path, output_path, to_fmt)
            
            if from_fmt in ('docx', 'doc') and to_fmt == 'pdf':
                return self._convert_word_pdf(input_path, output_path)
            
            if from_fmt in ('docx', 'doc') and to_fmt in ('txt', 'html'):
                return self._convert_word(input_path, output_path, to_fmt)
            
            if from_fmt in self.ARCHIVE_FORMATS:
                return self._convert_archive(input_path, output_path, from_fmt, to_fmt)
            
            if from_fmt == 'markdown' and to_fmt == 'html':
                return self.convert(input_path, 'markdown', 'html', output_path)
            
            if from_fmt in self.TEXT_FORMATS or from_fmt in self.DATA_FORMATS:
                with open(input_path, 'r', encoding='utf-8') as f:
                    data = f.read()
                return self.convert(data, from_fmt, to_fmt, output_path)
            
            if from_fmt in self.DATABASE_FORMATS:
                return self._convert_database(input_path, output_path, from_fmt, to_fmt)
            
            return ConversionResult(success=False, error=f"Unsupported: {from_fmt} → {to_fmt}")
        except Exception as e:
            return ConversionResult(success=False, error=str(e))
    
    def _convert_database(self, input_path: str, output_path: str, from_fmt: str, to_fmt: str) -> ConversionResult:
        try:
            import sqlite3
            
            conn = sqlite3.connect(input_path)
            cursor = conn.cursor()
            
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
            tables = cursor.fetchall()
            
            db_data = {'tables': {}}
            
            for (table_name,) in tables:
                cursor.execute(f"PRAGMA table_info({table_name})")
                columns = [row[1] for row in cursor.fetchall()]
                
                cursor.execute(f"SELECT * FROM {table_name}")
                rows = cursor.fetchall()
                
                db_data['tables'][table_name] = {
                    'columns': columns,
                    'data': [dict(zip(columns, row)) for row in rows]
                }
            
            conn.close()
            
            if to_fmt == 'json':
                result = json.dumps(db_data, indent=2, default=str)
            elif to_fmt == 'csv':
                all_data = []
                for tbl, data in db_data['tables'].items():
                    for row in data['data']:
                        row['_table'] = tbl
                        all_data.append(row)
                if all_data:
                    df = pd.DataFrame(all_data) if pd else None
                    if df is not None:
                        result = df.to_csv(index=False)
                    else:
                        result = self._to_csv({'data': all_data}, {})
                else:
                    result = ""
            elif to_fmt in ('mysql', 'postgresql', 'sqlite'):
                result = self._to_sql(db_data, {})
            else:
                return ConversionResult(success=False, error=f"Unsupported database conversion: {from_fmt} → {to_fmt}")
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(result)
            
            return ConversionResult(success=True, output_path=output_path, metadata={'format': to_fmt})
            
        except ImportError:
            return ConversionResult(success=False, error="sqlite3 (built-in) or pandas required for database conversion")
        except Exception as e:
            return ConversionResult(success=False, error=f"Database conversion failed: {e}")
            
        except Exception as e:
            return ConversionResult(success=False, error=str(e), details=traceback.format_exc())
    
    def _convert_image(self, input_path: str, output_path: str, from_fmt: str, to_fmt: str) -> ConversionResult:
        try:
            from PIL import Image
        except ImportError:
            return ConversionResult(success=False, error="Pillow required: pip install pillow")
        
        try:
            img = Image.open(input_path)
            
            if to_fmt in ('jpg', 'jpeg') and img.mode in ('RGBA', 'LA', 'P'):
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                img = background
            
            img.save(output_path, format=to_fmt.upper())
            return ConversionResult(success=True, output_path=output_path, metadata={'format': to_fmt})
        except Exception as e:
            return ConversionResult(success=False, error=f"Image conversion failed: {e}")
    
    def resize_image(self, input_path: str, output_path: str = None, width: int = None, height: int = None, scale: float = None, keep_aspect: bool = True) -> ConversionResult:
        try:
            from PIL import Image
        except ImportError:
            return ConversionResult(success=False, error="Pillow required: pip install pillow")
        
        try:
            img = Image.open(input_path)
            orig_w, orig_h = img.size
            
            if scale:
                width = int(orig_w * scale)
                height = int(orig_h * scale)
            elif width and height:
                if keep_aspect:
                    img.thumbnail((width, height), Image.LANCZOS)
                    width, height = img.size
                else:
                    pass
            elif width:
                height = int(orig_h * width / orig_w) if keep_aspect else width
            elif height:
                width = int(orig_w * height / orig_h) if keep_aspect else height
            
            if keep_aspect and width and height:
                img.thumbnail((width, height), Image.LANCZOS)
            elif width and height:
                img = img.resize((width, height), Image.LANCZOS)
            
            if output_path is None:
                output_path = input_path
            
            img.save(output_path)
            return ConversionResult(success=True, output_path=output_path, 
                                  metadata={'format': Path(output_path).suffix[1:], 'size': f"{img.size[0]}x{img.size[1]}"})
        except Exception as e:
            return ConversionResult(success=False, error=str(e))
    
    def _pdf_to_text(self, input_path: str, output_path: str) -> ConversionResult:
        try:
            import PyPDF2
            with open(input_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = '\n\n'.join(page.extract_text() or '' for page in reader.pages)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            return ConversionResult(success=True, output_path=output_path, metadata={'format': 'txt'})
        except ImportError:
            return ConversionResult(success=False, error="PyPDF2 required: pip install PyPDF2")
        except Exception as e:
            return ConversionResult(success=False, error=str(e))
    
    def _text_to_pdf(self, input_path: str, output_path: str) -> ConversionResult:
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            with open(input_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter
            text_object = c.beginText(50, height - 50)
            text_object.setFont("Helvetica", 10)
            
            for line in text.split('\n'):
                if line:
                    text_object.textLine(line)
            
            c.drawText(text_object)
            c.save()
            return ConversionResult(success=True, output_path=output_path, metadata={'format': 'pdf'})
        except ImportError:
            return ConversionResult(success=False, error="reportlab required: pip install reportlab")
        except Exception as e:
            return ConversionResult(success=False, error=str(e))
    
    def _convert_pdf_word(self, input_path: str, output_path: str, to_format: str) -> ConversionResult:
        temp_txt = tempfile.mktemp(suffix='.txt')
        result = self._pdf_to_text(input_path, temp_txt)
        if not result.success:
            return result
        
        try:
            from docx import Document
            doc = Document()
            
            with open(temp_txt, 'r', encoding='utf-8') as f:
                for para in f.read().split('\n\n'):
                    if para.strip():
                        doc.add_paragraph(para.strip())
            
            doc.save(output_path)
            os.unlink(temp_txt)
            return ConversionResult(success=True, output_path=output_path, metadata={'format': to_format})
        except ImportError:
            return ConversionResult(success=False, error="python-docx required: pip install python-docx")
        except Exception as e:
            return ConversionResult(success=False, error=str(e))
    
    def _convert_word_pdf(self, input_path: str, output_path: str) -> ConversionResult:
        return ConversionResult(success=False, error="Word to PDF requires LibreOffice or pandoc. Install: brew install libreoffice")
    
    def _convert_word(self, input_path: str, output_path: str, to_fmt: str) -> ConversionResult:
        try:
            from docx import Document
            doc = Document(input_path)
            
            if to_fmt == 'txt':
                text = '\n\n'.join(para.text for para in doc.paragraphs)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
            elif to_fmt == 'html':
                html = ['<!DOCTYPE html><html><head><meta charset="UTF-8"></head><body>']
                for para in doc.paragraphs:
                    if para.text.strip():
                        html.append(f'<p>{para.text}</p>')
                html.append('</body></html>')
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write('\n'.join(html))
            
            return ConversionResult(success=True, output_path=output_path, metadata={'format': to_fmt})
        except ImportError:
            return ConversionResult(success=False, error="python-docx required: pip install python-docx")
        except Exception as e:
            return ConversionResult(success=False, error=str(e))
    
    def _convert_archive(self, input_path: str, output_path: str, from_fmt: str, to_fmt: str) -> ConversionResult:
        try:
            if from_fmt == 'zip' and to_fmt == 'tar':
                with zipfile.ZipFile(input_path, 'r') as z:
                    with tarfile.open(output_path, 'w') as t:
                        for name in z.namelist():
                            t.addfile(z.open(name), arcname=name)
                return ConversionResult(success=True, output_path=output_path)
            
            if from_fmt == 'tar' and to_fmt == 'zip':
                with tarfile.open(input_path, 'r') as t:
                    with zipfile.ZipFile(output_path, 'w') as z:
                        for member in t.getmembers():
                            f = t.extractfile(member)
                            if f:
                                z.writestr(member.name, f.read())
                return ConversionResult(success=True, output_path=output_path)
            
            return ConversionResult(success=False, error=f"Archive conversion {from_fmt} → {to_fmt} not supported")
        except Exception as e:
            return ConversionResult(success=False, error=str(e))
    
    def batch_convert(self, input_files: List[str], output_dir: str = None, to_format: str = None) -> List[ConversionResult]:
        results = []
        
        if output_dir:
            Path(output_dir).mkdir(parents=True, exist_ok=True)
        
        for input_file in input_files:
            ip = Path(input_file)
            output_format = to_format or ip.suffix[1:].lower()
            op = str(Path(output_dir) / f"{ip.stem}.{output_format}") if output_dir else None
            result = self.convert_file(input_file, op)
            results.append(result)
        
        return results
    
    def get_supported_formats(self) -> Dict[str, List[str]]:
        return {
            'text': sorted(self.TEXT_FORMATS),
            'data': sorted(self.DATA_FORMATS),
            'document': sorted(self.DOCUMENT_FORMATS),
            'image': sorted(self.IMAGE_FORMATS),
            'archive': sorted(self.ARCHIVE_FORMATS),
            'code': sorted(self.CODE_FORMATS),
        }
    
    def all_formats(self) -> List[str]:
        return sorted(self.ALL_FORMATS)


def convert(data: Any, from_fmt: str = None, to_fmt: str = None, output: str = None, **kwargs) -> Any:
    """Quick conversion function"""
    converter = UniversalConverter()
    result = converter.convert(data, from_fmt, to_fmt, output, **kwargs)
    if result.success:
        return result.data or result.output_path
    raise ConversionError(result.error or "Conversion failed")


def convert_file(input_path: str, output_path: str = None, from_fmt: str = None, to_fmt: str = None) -> str:
    """Convert file with format detection"""
    converter = UniversalConverter()
    result = converter.convert_file(input_path, output_path, from_fmt, to_fmt)
    if result.success:
        return result.output_path or result.data
    raise ConversionError(result.error)


def resize_image(input_path: str, output_path: str = None, width: int = None, height: int = None, scale: float = None) -> str:
    """Resize image"""
    converter = UniversalConverter()
    result = converter.resize_image(input_path, output_path, width, height, scale)
    if result.success:
        return result.output_path
    raise ConversionError(result.error)


def batch_convert(files: List[str], output_dir: str = None, to_format: str = 'json') -> List[Dict]:
    """Batch convert files"""
    converter = UniversalConverter()
    results = converter.batch_convert(files, output_dir, to_format)
    return [{'file': f, 'success': r.success, 'output': r.output_path, 'error': r.error} 
            for f, r in zip(files, results)]


def get_formats() -> Dict[str, List[str]]:
    return UniversalConverter().get_supported_formats()


def all_supported_formats() -> List[str]:
    return UniversalConverter().all_formats()


def file_info(file_path: str) -> Dict:
    """Get file information"""
    p = Path(file_path)
    if not p.exists():
        return {'valid': False, 'error': 'File not found'}
    
    stat = p.stat()
    return {
        'valid': True,
        'name': p.name,
        'size': stat.st_size,
        'size_human': f"{stat.st_size / 1024:.2f} KB",
        'format': p.suffix[1:].lower() or None,
        'mime': mimetypes.guess_type(str(p))[0],
        'modified': datetime.fromtimestamp(stat.st_mtime).isoformat(),
        'created': datetime.fromtimestamp(stat.st_ctime).isoformat(),
        'permissions': oct(stat.st_mode)[-3:],
    }


def hash_file(file_path: str, algorithm: str = 'sha256') -> str:
    """Hash a file"""
    h = hashlib.new(algorithm)
    with open(file_path, 'rb') as f:
        for chunk in iter(lambda: f.read(8192), b''):
            h.update(chunk)
    return h.hexdigest()


def validate_file(file_path: str) -> Dict:
    """Validate file format"""
    p = Path(file_path)
    if not p.exists():
        return {'valid': False, 'error': 'File not found'}
    
    fmt = p.suffix[1:].lower()
    supported = UniversalConverter().all_formats()
    
    return {
        'valid': fmt in supported,
        'format': fmt,
        'supported': fmt in supported,
        'category': _get_format_category(fmt),
    }


def _get_format_category(fmt: str) -> str:
    c = UniversalConverter()
    if fmt in c.TEXT_FORMATS: return 'text'
    if fmt in c.DATA_FORMATS: return 'data'
    if fmt in c.DOCUMENT_FORMATS: return 'document'
    if fmt in c.IMAGE_FORMATS: return 'image'
    if fmt in c.ARCHIVE_FORMATS: return 'archive'
    if fmt in c.CODE_FORMATS: return 'code'
    if fmt in c.DATABASE_FORMATS: return 'database'
    return 'unknown'


# ============ ENCODING CONVERSIONS ============

def encode_base64(data: Union[str, bytes]) -> str:
    """Encode string or bytes to Base64"""
    if isinstance(data, str):
        data = data.encode('utf-8')
    return base64.b64encode(data).decode('utf-8')


def decode_base64(data: str) -> str:
    """Decode Base64 to string"""
    return base64.b64decode(data).decode('utf-8')


def encode_hex(data: Union[str, bytes]) -> str:
    """Encode string or bytes to Hex"""
    if isinstance(data, str):
        data = data.encode('utf-8')
    return data.hex()


def decode_hex(data: str) -> str:
    """Decode Hex to string"""
    return bytes.fromhex(data).decode('utf-8')


def encode_url(data: str) -> str:
    """URL encode a string"""
    from urllib.parse import quote
    return quote(data)


def decode_url(data: str) -> str:
    """URL decode a string"""
    from urllib.parse import unquote
    return unquote(data)


def encode_html(data: str) -> str:
    """HTML encode a string"""
    return (data.replace('&', '&amp;')
                .replace('<', '&lt;')
                .replace('>', '&gt;')
                .replace('"', '&quot;')
                .replace("'", '&#39;'))


def decode_html(data: str) -> str:
    """HTML decode a string"""
    import html
    return html.unescape(data)


# ============ COLOR CONVERSIONS ============

def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """Convert HEX color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join(c*2 for c in hex_color)
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def rgb_to_hex(r: int, g: int, b: int) -> str:
    """Convert RGB to HEX color"""
    return f'#{r:02x}{g:02x}{b:02x}'


def rgb_to_hsl(r: int, g: int, b: int) -> Tuple[float, float, float]:
    """Convert RGB to HSL"""
    r, g, b = r/255, g/255, b/255
    max_c = max(r, g, b)
    min_c = min(r, g, b)
    l = (max_c + min_c) / 2
    
    if max_c == min_c:
        h = s = 0
    else:
        d = max_c - min_c
        s = d / (2 - max_c - min_c) if l > 0.5 else d / (max_c + min_c)
        if max_c == r:
            h = ((g - b) / d + (6 if g < b else 0)) / 6
        elif max_c == g:
            h = ((b - r) / d + 2) / 6
        else:
            h = ((r - g) / d + 4) / 6
    return round(h * 360, 1), round(s * 100, 1), round(l * 100, 1)


def hsl_to_rgb(h: float, s: float, l: float) -> Tuple[int, int, int]:
    """Convert HSL to RGB"""
    h, s, l = h/360, s/100, l/100
    
    if s == 0:
        val = round(l * 255)
        return (val, val, val)
    
    def hue_to_rgb(p, q, t):
        if t < 0: t += 1
        if t > 1: t -= 1
        if t < 1/6: return p + (q - p) * 6 * t
        if t < 1/2: return q
        if t < 2/3: return p + (q - p) * (2/3 - t) * 6
        return p
    
    q = l * (1 + s) if l < 0.5 else l + s - l * s
    p = 2 * l - q
    r = round(hue_to_rgb(p, q, h + 1/3) * 255)
    g = round(hue_to_rgb(p, q, h) * 255)
    b = round(hue_to_rgb(p, q, h - 1/3) * 255)
    return (r, g, b)


def convert_color(color: str, to_format: str) -> Any:
    """Convert between color formats (hex, rgb, hsl)"""
    color = color.strip().lower()
    
    if to_format == 'hex':
        if color.startswith('rgb'):
            match = re.match(r'rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color)
            if match:
                return rgb_to_hex(int(match[1]), int(match[2]), int(match[3]))
        elif color.startswith('hsl'):
            match = re.match(r'hsl\s*\(\s*(\d+)\s*,\s*(\d+)%?\s*,\s*(\d+)%?\s*\)', color)
            if match:
                rgb = hsl_to_rgb(int(match[1]), int(match[2]), int(match[3]))
                return rgb_to_hex(*rgb)
        return color if color.startswith('#') else color
    
    if to_format == 'rgb':
        if color.startswith('#'):
            rgb = hex_to_rgb(color)
            return f"rgb({rgb[0]}, {rgb[1]}, {rgb[2]})"
        elif color.startswith('hsl'):
            match = re.match(r'hsl\s*\(\s*(\d+)\s*,\s*(\d+)%?\s*,\s*(\d+)%?\s*\)', color)
            if match:
                rgb = hsl_to_rgb(int(match[1]), int(match[2]), int(match[3]))
                return f"rgb({rgb[0]}, {rgb[1]}, {rgb[2]})"
    
    if to_format == 'hsl':
        if color.startswith('#'):
            rgb = hex_to_rgb(color)
            hsl = rgb_to_hsl(*rgb)
            return f"hsl({hsl[0]}, {hsl[1]}%, {hsl[2]}%)"
        elif color.startswith('rgb'):
            match = re.match(r'rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', color)
            if match:
                hsl = rgb_to_hsl(int(match[1]), int(match[2]), int(match[3]))
                return f"hsl({hsl[0]}, {hsl[1]}%, {hsl[2]}%)"
    
    return None


# ============ DATE/TIME CONVERSIONS ============

def timestamp_to_iso(timestamp: Union[int, float], utc: bool = False) -> str:
    """Convert Unix timestamp to ISO format"""
    dt = datetime.fromtimestamp(timestamp, datetime.timezone.utc if utc else None)
    return dt.isoformat()


def iso_to_timestamp(iso_str: str) -> float:
    """Convert ISO string to Unix timestamp"""
    dt = datetime.fromisoformat(iso_str.replace('Z', '+00:00'))
    return dt.timestamp()


def timestamp_to_human(timestamp: Union[int, float], format_str: str = '%Y-%m-%d %H:%M:%S') -> str:
    """Convert Unix timestamp to human readable format"""
    dt = datetime.fromtimestamp(timestamp)
    return dt.strftime(format_str)


def human_to_timestamp(date_str: str, format_str: str = '%Y-%m-%d %H:%M:%S') -> float:
    """Convert human readable date to Unix timestamp"""
    dt = datetime.strptime(date_str, format_str)
    return dt.timestamp()


def convert_date(date_str: str, from_format: str, to_format: str) -> str:
    """Convert between date formats"""
    dt = datetime.strptime(date_str, from_format)
    return dt.strftime(to_format)


def get_current_timestamp() -> int:
    """Get current Unix timestamp"""
    return int(datetime.now().timestamp())


def get_current_iso() -> str:
    """Get current ISO timestamp"""
    return datetime.now().isoformat()


# ============ CASE CONVERSIONS ============

def to_camel_case(text: str) -> str:
    """Convert to camelCase"""
    text = re.sub(r'[^a-zA-Z0-9\s]', ' ', text)
    words = text.split()
    if not words: return text
    return words[0].lower() + ''.join(word.capitalize() for word in words[1:])


def to_pascal_case(text: str) -> str:
    """Convert to PascalCase"""
    text = re.sub(r'[^a-zA-Z0-9\s]', ' ', text)
    words = text.split()
    return ''.join(word.capitalize() for word in words)


def to_snake_case(text: str) -> str:
    """Convert to snake_case"""
    text = re.sub(r'[^a-zA-Z0-9\s]', ' ', text)
    text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
    return '_'.join(word.lower() for word in text.split())


def to_kebab_case(text: str) -> str:
    """Convert to kebab-case"""
    text = re.sub(r'[^a-zA-Z0-9\s]', ' ', text)
    text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)
    return '-'.join(word.lower() for word in text.split())


def to_screaming_snake_case(text: str) -> str:
    """Convert to SCREAMING_SNAKE_CASE"""
    return to_snake_case(text).upper()


def to_sentence_case(text: str) -> str:
    """Convert to Sentence case"""
    text = text.lower().strip()
    return text[0].upper() + text[1:] if text else text


def convert_case(text: str, to_case: str) -> str:
    """Convert between case formats"""
    cases = {
        'camel': to_camel_case,
        'pascal': to_pascal_case,
        'snake': to_snake_case,
        'kebab': to_kebab_case,
        'screaming': to_screaming_snake_case,
        'sentence': to_sentence_case,
    }
    converter = cases.get(to_case.lower())
    return converter(text) if converter else text


# ============ UNIT CONVERSIONS ============

UNITS = {
    'length': {
        'm': 1.0, 'meter': 1.0, 'meters': 1.0,
        'km': 1000.0, 'kilometer': 1000.0,
        'cm': 0.01, 'centimeter': 0.01,
        'mm': 0.001, 'millimeter': 0.001,
        'mi': 1609.344, 'mile': 1609.344,
        'yd': 0.9144, 'yard': 0.9144,
        'ft': 0.3048, 'foot': 0.3048, 'feet': 0.3048,
        'in': 0.0254, 'inch': 0.0254,
    },
    'weight': {
        'kg': 1.0, 'kilogram': 1.0, 'kilograms': 1.0,
        'g': 0.001, 'gram': 0.001, 'grams': 0.001,
        'mg': 0.000001, 'milligram': 0.000001,
        'lb': 0.453592, 'pound': 0.453592, 'pounds': 0.453592,
        'oz': 0.0283495, 'ounce': 0.0283495,
    },
    'temperature': {
        'c': 1.0, 'celsius': 1.0,
        'f': 1.0, 'fahrenheit': 1.0,
        'k': 1.0, 'kelvin': 1.0,
    },
    'volume': {
        'l': 1.0, 'liter': 1.0, 'liters': 1.0,
        'ml': 0.001, 'milliliter': 0.001,
        'gal': 3.78541, 'gallon': 3.78541,
        'qt': 0.946353, 'quart': 0.946353,
        'pt': 0.473176, 'pint': 0.473176,
        'cup': 0.236588, 'cup': 0.236588,
        'fl oz': 0.0295735, 'fluid ounce': 0.0295735,
    },
    'data': {
        'b': 1.0, 'byte': 1.0, 'bytes': 1.0,
        'kb': 1024.0, 'kilobyte': 1024.0,
        'mb': 1048576.0, 'megabyte': 1048576.0,
        'gb': 1073741824.0, 'gigabyte': 1073741824.0,
        'tb': 1099511627776.0, 'terabyte': 1099511627776.0,
    },
}


def _get_unit_type(unit: str) -> str:
    """Detect unit type"""
    unit = unit.lower()
    for utype, units in UNITS.items():
        if unit in units:
            return utype
    return None


def convert_unit(value: float, from_unit: str, to_unit: str) -> Optional[float]:
    """Convert between units"""
    from_unit = from_unit.lower()
    to_unit = to_unit.lower()
    
    unit_type = _get_unit_type(from_unit)
    if not unit_type or _get_unit_type(to_unit) != unit_type:
        return None
    
    if unit_type == 'temperature':
        if from_unit in ('c', 'celsius'):
            c = value
        elif from_unit in ('f', 'fahrenheit'):
            c = (value - 32) * 5/9
        else:
            c = value - 273.15
        
        if to_unit in ('c', 'celsius'):
            return c
        elif to_unit in ('f', 'fahrenheit'):
            return c * 9/5 + 32
        else:
            return c + 273.15
    
    base_value = value * UNITS[unit_type][from_unit]
    return base_value / UNITS[unit_type][to_unit]


# ============ HASH FUNCTIONS ============

def hash_string(text: str, algorithm: str = 'sha256') -> str:
    """Hash a string"""
    h = hashlib.new(algorithm)
    h.update(text.encode('utf-8'))
    return h.hexdigest()


def hash_file_md5(file_path: str) -> str:
    """MD5 hash of file"""
    return hash_file(file_path, 'md5')


def hash_file_sha1(file_path: str) -> str:
    """SHA1 hash of file"""
    return hash_file(file_path, 'sha1')


def hash_file_sha256(file_path: str) -> str:
    """SHA256 hash of file"""
    return hash_file(file_path, 'sha256')


def verify_hash(file_path: str, expected_hash: str, algorithm: str = 'sha256') -> bool:
    """Verify file hash"""
    actual = hash_file(file_path, algorithm)
    return actual.lower() == expected_hash.lower()


# ============ SERIALIZATION ============

def to_pickle(obj: Any) -> bytes:
    """Serialize to Pickle"""
    import pickle
    return pickle.dumps(obj)


def from_pickle(data: bytes) -> Any:
    """Deserialize from Pickle"""
    import pickle
    return pickle.loads(data)


def to_msgpack(obj: Any) -> bytes:
    """Serialize to MessagePack"""
    try:
        import msgpack
        return msgpack.packb(obj, use_bin_type=True)
    except ImportError:
        raise ConversionError("msgpack not installed: pip install msgpack")


def from_msgpack(data: bytes) -> Any:
    """Deserialize from MessagePack"""
    try:
        import msgpack
        return msgpack.unpackb(data, raw=False)
    except ImportError:
        raise ConversionError("msgpack not installed: pip install msgpack")


def to_toml(obj: Dict) -> str:
    """Serialize to TOML"""
    try:
        import tomli
        return tomli.dumps(obj)
    except ImportError:
        return json.dumps(obj, indent=2)


def from_toml(data: str) -> Dict:
    """Deserialize from TOML"""
    try:
        import tomli
        return tomli.loads(data)
    except ImportError:
        return json.loads(data)


# ============ MIME TYPES ============

def get_mime_type(file_path: str) -> Optional[str]:
    """Get MIME type from file extension"""
    mime, _ = mimetypes.guess_type(file_path)
    return mime


def get_extension(mime_type: str) -> Optional[str]:
    """Get file extension from MIME type"""
    ext = mimetypes.guess_extension(mime_type)
    return ext


def is_text_mime(mime_type: str) -> bool:
    """Check if MIME type is text"""
    return mime_type and mime_type.startswith('text/')


def is_image_mime(mime_type: str) -> bool:
    """Check if MIME type is image"""
    return mime_type and mime_type.startswith('image/')


# ============ FILETYPE DETECTION ============

def detect_file_type(file_path: str) -> Dict:
    """Detect file type using magic bytes"""
    p = Path(file_path)
    if not p.exists():
        return {'error': 'File not found'}
    
    signatures = {
        b'\x89PNG': ('png', 'image/png'),
        b'\xff\xd8\xff': ('jpg', 'image/jpeg'),
        b'GIF87a': ('gif', 'image/gif'),
        b'GIF89a': ('gif', 'image/gif'),
        b'BM': ('bmp', 'image/bmp'),
        b'PK\x03\x04': ('zip', 'application/zip'),
        b'\x1f\x8b': ('gz', 'application/gzip'),
        b'SQLite format 3': ('sqlite', 'application/x-sqlite3'),
        b'%PDF': ('pdf', 'application/pdf'),
        b'\xca\xfe\xba\xbe': ('class', 'application/java'),
        b'\x50\x4b': ('docx', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'),
    }
    
    with open(file_path, 'rb') as f:
        header = f.read(16)
    
    for sig, (ext, mime) in signatures.items():
        if header.startswith(sig):
            return {'extension': ext, 'mime': mime, 'format': 'binary'}
    
    return {'extension': p.suffix[1:] or None, 'mime': get_mime_type(file_path), 'format': 'text' if is_text_mime(get_mime_type(file_path)) else 'unknown'}


# ============ MARKDOWN/HTML/AAGENT/CLAUDE ============

def markdown_to_html(md: str) -> str:
    """Convert Markdown to HTML"""
    import re
    
    html = md
    
    headers = re.findall(r'^(#{1,6})\s+(.+)$', md, re.MULTILINE)
    for hashes, text in headers:
        level = len(hashes)
        html = html.replace(f"{hashes} {text}", f"<h{level}>{text}</h{level}>")
    
    bold = re.findall(r'\*\*(.+?)\*\*', html)
    for text in bold:
        html = html.replace(f"**{text}**", f"<strong>{text}</strong>")
    
    italic = re.findall(r'\*(.+?)\*', html)
    for text in italic:
        html = html.replace(f"*{text}*", f"<em>{text}</em>")
    
    code_blocks = re.findall(r'```(\w*)\n([\s\S]*?)```', html)
    for lang, code in code_blocks:
        html = html.replace(f"```{lang}\n{code}```", f'<pre><code class="language-{lang}">{code}</code></pre>')
    
    inline_code = re.findall(r'`(.+?)`', html)
    for code in inline_code:
        html = html.replace(f"`{code}`", f"<code>{code}</code>")
    
    links = re.findall(r'\[(.+?)\]\((.+?)\)', html)
    for text, url in links:
        html = html.replace(f"[{text}]({url})", f'<a href="{url}">{text}</a>')
    
    lines = html.split('\n')
    result = ['<html><head><meta charset="UTF-8"></head><body>']
    in_list = False
    
    for line in lines:
        line = line.strip()
        if line.startswith('<h') or line.startswith('<pre'):
            result.append(line)
            in_list = False
        elif line.startswith('- ') or line.startswith('* '):
            if not in_list:
                result.append('<ul>')
                in_list = True
            result.append(f"<li>{line[2:]}</li>")
        elif line:
            if in_list:
                result.append('</ul>')
                in_list = False
            result.append(f"<p>{line}</p>")
    
    if in_list:
        result.append('</ul>')
    result.append('</body></html>')
    
    return '\n'.join(result)


def html_to_markdown(html: str) -> str:
    """Convert HTML to Markdown"""
    import re
    
    md = html
    
    h_tags = re.findall(r'<h(\d)>(.+?)</h\d>', html)
    for level, text in h_tags:
        md = md.replace(f"<h{level}>{text}</h{level}>", f"{'#' * int(level)} {text}")
    
    md = re.sub(r'<strong>(.+?)</strong>', r'**\1**', md)
    md = re.sub(r'<b>(.+?)</b>', r'**\1**', md)
    md = re.sub(r'<em>(.+?)</em>', r'*\1*', md)
    md = re.sub(r'<i>(.+?)</i>', r'*\1*', md)
    md = re.sub(r'<code>(.+?)</code>', r'`\1`', md)
    md = re.sub(r'<a href="(.+?)">(.+?)</a>', r'[\2](\1)', md)
    md = re.sub(r'<li>(.+?)</li>', r'- \1', md)
    md = re.sub(r'<p>(.+?)</p>', r'\1\n', md)
    md = re.sub(r'<br\s*/?>', r'\n', md)
    md = re.sub(r'<[^>]+>', '', md)
    md = re.sub(r'\n{3,}', '\n\n', md)
    
    return md.strip()


def to_agent_format(text: str) -> str:
    """Convert to agent/Claude format context"""
    return f"""# Context

{text}

---

## Guidelines
- Think silently when needed
- Use tools for file operations
- Ask for clarification when needed
- Provide concise responses
"""


def to_claude_format(text: str) -> str:
    """Convert to Claude Code context format"""
    return f"""# Claude Code Context

## Project Overview
This document contains context for Claude Code AI assistant.

## Guidelines
- Think silently when needed
- Use tools for file operations
- Ask for clarification when needed
- Provide concise responses

---

{text}
"""


# ============ SPREADSHEET CONVERSIONS ============

def csv_to_excel(csv_path: str, excel_path: str = None) -> str:
    """Convert CSV to Excel (.xlsx)"""
    try:
        import pandas as pd
        df = pd.read_csv(csv_path)
        if excel_path is None:
            excel_path = csv_path.rsplit('.', 1)[0] + '.xlsx'
        df.to_excel(excel_path, index=False)
        return excel_path
    except ImportError:
        raise ConversionError("pandas and openpyxl required: pip install pandas openpyxl")


def excel_to_csv(excel_path: str, csv_path: str = None, sheet_name: int = 0) -> str:
    """Convert Excel to CSV"""
    try:
        import pandas as pd
        df = pd.read_excel(excel_path, sheet_name=sheet_name)
        if csv_path is None:
            csv_path = excel_path.rsplit('.', 1)[0] + '.csv'
        df.to_csv(csv_path, index=False)
        return csv_path
    except ImportError:
        raise ConversionError("pandas and openpyxl required: pip install pandas openpyxl")


def xls_to_xlsx(xls_path: str, xlsx_path: str = None) -> str:
    """Convert old Excel (.xls) to new format (.xlsx)"""
    try:
        import pandas as pd
        df = pd.read_excel(xls_path, engine='xlrd')
        if xlsx_path is None:
            xlsx_path = xls_path.rsplit('.', 1)[0] + '.xlsx'
        df.to_excel(xlsx_path, index=False)
        return xlsx_path
    except ImportError:
        raise ConversionError("pandas, xlrd, openpyxl required: pip install pandas xlrd openpyxl")


def excel_to_json(excel_path: str, json_path: str = None, sheet_name: int = 0) -> str:
    """Convert Excel to JSON"""
    try:
        import pandas as pd
        df = pd.read_excel(excel_path, sheet_name=sheet_name)
        if json_path is None:
            json_path = excel_path.rsplit('.', 1)[0] + '.json'
        df.to_json(json_path, orient='records', indent=2)
        return json_path
    except ImportError:
        raise ConversionError("pandas and openpyxl required: pip install pandas openpyxl")


def json_to_excel(json_path: str, excel_path: str = None) -> str:
    """Convert JSON to Excel"""
    try:
        import pandas as pd
        df = pd.read_json(json_path)
        if excel_path is None:
            excel_path = json_path.rsplit('.', 1)[0] + '.xlsx'
        df.to_excel(excel_path, index=False)
        return excel_path
    except ImportError:
        raise ConversionError("pandas and openpyxl required: pip install pandas openpyxl")


# ============ PRESENTATION CONVERSIONS ============

def text_to_pptx(text_path: str, pptx_path: str = None, title: str = "Presentation") -> str:
    """Convert text file to PowerPoint"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        with open(text_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        content_slide.shapes.title.text = "Content"
        content_slide.placeholders[1].text = content
        
        if pptx_path is None:
            pptx_path = text_path.rsplit('.', 1)[0] + '.pptx'
        
        prs.save(pptx_path)
        return pptx_path
    except ImportError:
        raise ConversionError("python-pptx required: pip install python-pptx")


def pptx_to_pdf(pptx_path: str, pdf_path: str = None) -> str:
    """Convert PowerPoint to PDF (requires LibreOffice)"""
    import subprocess
    import os
    
    if pdf_path is None:
        pdf_path = pptx_path.rsplit('.', 1)[0] + '.pdf'
    
    cmd = ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', 
           os.path.dirname(pdf_path), pptx_path]
    
    try:
        subprocess.run(cmd, check=True, capture_output=True)
        return pdf_path
    except FileNotFoundError:
        raise ConversionError("LibreOffice required: brew install libreoffice")


def pptx_to_text(pptx_path: str, txt_path: str = None) -> str:
    """Convert PowerPoint to text"""
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        text_lines = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_lines.append(shape.text)
        
        text_content = '\n\n'.join(text_lines)
        
        if txt_path is None:
            txt_path = pptx_path.rsplit('.', 1)[0] + '.txt'
        
        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write(text_content)
        
        return txt_path
    except ImportError:
        raise ConversionError("python-pptx required: pip install python-pptx")


# ============ MARKDOWN/HTML TO PDF ============

def markdown_to_pdf(md_path: str, pdf_path: str = None) -> str:
    """Convert Markdown to PDF"""
    try:
        import markdown
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from reportlab.lib.styles import getSampleStyleSheet
        
        with open(md_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        html = markdown.markdown(md_content)
        
        if pdf_path is None:
            pdf_path = md_path.rsplit('.', 1)[0] + '.pdf'
        
        c = canvas.Canvas(pdf_path, pagesize=letter)
        width, height = letter
        
        styles = getSampleStyleSheet()
        style = styles['Normal']
        
        text_object = c.beginText(50, height - 50)
        text_object.setFont("Helvetica", 10)
        
        for line in html.replace('<br>', '\n').replace('<p>', '').replace('</p>', '\n').split('\n'):
            if line.strip():
                text_object.textLine(line[:100])
        
        c.drawText(text_object)
        c.save()
        
        return pdf_path
    except ImportError:
        raise ConversionError("markdown and reportlab required: pip install markdown reportlab")


def html_to_pdf(html_path: str, pdf_path: str = None) -> str:
    """Convert HTML to PDF"""
    try:
        from weasyprint import HTML
    except ImportError:
        raise ConversionError("weasyprint required: pip install weasyprint")
    
    if pdf_path is None:
        pdf_path = html_path.rsplit('.', 1)[0] + '.pdf'
    
    HTML(filename=html_path).write_pdf(pdf_path)
    return pdf_path


# ============ RTF CONVERSIONS ============

def rtf_to_text(rtf_path: str, txt_path: str = None) -> str:
    """Convert RTF to plain text"""
    import re
    
    with open(rtf_path, 'r', encoding='utf-8', errors='ignore') as f:
        rtf_content = f.read()
    
    text = rtf_content
    text = re.sub(r'\\[a-z]+\d*\s?', '', text)
    text = re.sub(r'[{}\\]', '', text)
    text = text.strip()
    
    if txt_path is None:
        txt_path = rtf_path.rsplit('.', 1)[0] + '.txt'
    
    with open(txt_path, 'w', encoding='utf-8') as f:
        f.write(text)
    
    return txt_path


def text_to_rtf(txt_path: str, rtf_path: str = None) -> str:
    """Convert plain text to RTF"""
    with open(txt_path, 'r', encoding='utf-8') as f:
        text = f.read()
    
    rtf_content = r"{\rtf1\ansi\deff0 " + text.replace('\n', r'\par ').replace('\\', r'\\') + "}"
    
    if rtf_path is None:
        rtf_path = txt_path.rsplit('.', 1)[0] + '.rtf'
    
    with open(rtf_path, 'w', encoding='utf-8') as f:
        f.write(rtf_content)
    
    return rtf_path


# ============ ODT CONVERSIONS ============

def odt_to_text(odt_path: str, txt_path: str = None) -> str:
    """Convert ODT to plain text"""
    try:
        from odf import text, teletype
        from odf.opendocument import load
    
        doc = load(odt_path)
        allparas = doc.getElementsByType(text.P)
        content = '\n'.join(teletype.extractText(p) for p in allparas)
        
        if txt_path is None:
            txt_path = odt_path.rsplit('.', 1)[0] + '.txt'
        
        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return txt_path
    except ImportError:
        raise ConversionError("odfpy required: pip install odfpy")


def text_to_odt(txt_path: str, odt_path: str = None) -> str:
    """Convert plain text to ODT"""
    try:
        from odf.opendocument import Text
        from odf.text import P
        from odf import style
        from odf.style import Style, TextProperties, ParagraphProperties
        
        with open(txt_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        doc = Text()
        
        bold_style = Style(name="Bold", family="text")
        bold_style.addElement(TextProperties(fontweight="bold"))
        doc.styles.addElement(bold_style)
        
        for line in content.split('\n'):
            doc.addElement(P(text=line))
        
        if odt_path is None:
            odt_path = txt_path.rsplit('.', 1)[0] + '.odt'
        
        doc.save(odt_path)
        return odt_path
    except ImportError:
        raise ConversionError("odfpy required: pip install odfpy")


# ============ SVG CONVERSIONS ============

def svg_to_png(svg_path: str, png_path: str = None, width: int = None, height: int = None) -> str:
    """Convert SVG to PNG"""
    try:
        import cairosvg
    except ImportError:
        raise ConversionError("cairosvg required: pip install cairosvg")
    
    if png_path is None:
        png_path = svg_path.rsplit('.', 1)[0] + '.png'
    
    cairosvg.svg2png(url=svg_path, write_to=png_path, output_width=width, output_height=height)
    return png_path


def svg_to_pdf(svg_path: str, pdf_path: str = None) -> str:
    """Convert SVG to PDF"""
    try:
        import cairosvg
    except ImportError:
        raise ConversionError("cairosvg required: pip install cairosvg")
    
    if pdf_path is None:
        pdf_path = svg_path.rsplit('.', 1)[0] + '.pdf'
    
    cairosvg.svg2pdf(url=svg_path, write_to=pdf_path)
    return pdf_path


# ============ AUDIO CONVERSIONS ============

def mp3_to_wav(mp3_path: str, wav_path: str = None) -> str:
    """Convert MP3 to WAV"""
    try:
        from pydub import AudioSegment
    except ImportError:
        raise ConversionError("pydub required: pip install pydub")
    
    sound = AudioSegment.from_mp3(mp3_path)
    
    if wav_path is None:
        wav_path = mp3_path.rsplit('.', 1)[0] + '.wav'
    
    sound.export(wav_path, format="wav")
    return wav_path


def wav_to_mp3(wav_path: str, mp3_path: str = None, bitrate: str = "192k") -> str:
    """Convert WAV to MP3"""
    try:
        from pydub import AudioSegment
    except ImportError:
        raise ConversionError("pydub required: pip install pydub")
    
    sound = AudioSegment.from_wav(wav_path)
    
    if mp3_path is None:
        mp3_path = wav_path.rsplit('.', 1)[0] + '.mp3'
    
    sound.export(mp3_path, format="mp3", bitrate=bitrate)
    return mp3_path


def audio_to_text(audio_path: str, language: str = "en") -> str:
    """Convert audio to text (Speech-to-Text)"""
    try:
        import speech_recognition as sr
    except ImportError:
        raise ConversionError("SpeechRecognition required: pip install SpeechRecognition")
    
    recognizer = sr.Recognizer()
    
    if audio_path.endswith('.mp3'):
        import io
        from pydub import AudioSegment
        sound = AudioSegment.from_mp3(audio_path)
        audio_path = io.BytesIO()
        sound.export(audio_path, format="wav")
        audio_path.seek(0)
        with sr.AudioFile(audio_path) as source:
            audio_data = recognizer.record(source)
    else:
        with sr.AudioFile(audio_path) as source:
            audio_data = recognizer.record(source)
    
    text = recognizer.recognize_google(audio_data, language=language)
    return text


# ============ VIDEO CONVERSIONS ============

def video_to_gif(video_path: str, gif_path: str = None, fps: int = 10, width: int = None) -> str:
    """Convert video to GIF"""
    try:
        from moviepy.editor import VideoFileClip
    except ImportError:
        raise ConversionError("moviepy required: pip install moviepy")
    
    clip = VideoFileClip(video_path)
    
    if width:
        clip = clip.resize(width=width)
    
    if gif_path is None:
        gif_path = video_path.rsplit('.', 1)[0] + '.gif'
    
    clip.write_gif(gif_path, fps=fps)
    clip.close()
    
    return gif_path


def video_to_audio(video_path: str, audio_path: str = None, format: str = "mp3") -> str:
    """Extract audio from video"""
    try:
        from moviepy.editor import VideoFileClip
    except ImportError:
        raise ConversionError("moviepy required: pip install moviepy")
    
    clip = VideoFileClip(video_path)
    audio = clip.audio
    
    if audio_path is None:
        audio_path = video_path.rsplit('.', 1)[0] + '.' + format
    
    audio.write_audiofile(audio_path)
    clip.close()
    
    return audio_path


def mp4_to_avi(mp4_path: str, avi_path: str = None) -> str:
    """Convert MP4 to AVI"""
    try:
        from moviepy.editor import VideoFileClip
    except ImportError:
        raise ConversionError("moviepy required: pip install moviepy")
    
    clip = VideoFileClip(mp4_path)
    
    if avi_path is None:
        avi_path = mp4_path.rsplit('.', 1)[0] + '.avi'
    
    clip.write_videofile(avi_path, codec='rawvideo', audio_codec='pcm_s16le')
    clip.close()
    
    return avi_path


def avi_to_mp4(avi_path: str, mp4_path: str = None) -> str:
    """Convert AVI to MP4"""
    try:
        from moviepy.editor import VideoFileClip
    except ImportError:
        raise ConversionError("moviepy required: pip install moviepy")
    
    clip = VideoFileClip(avi_path)
    
    if mp4_path is None:
        mp4_path = avi_path.rsplit('.', 1)[0] + '.mp4'
    
    clip.write_videofile(mp4_path, codec='libx264')
    clip.close()
    
    return mp4_path


# ============ RAR ARCHIVE ============

def extract_rar(rar_path: str, extract_to: str = None) -> str:
    """Extract RAR archive"""
    try:
        import rarfile
    except ImportError:
        raise ConversionError("rarfile required: pip install rarfile")
    
    if extract_to is None:
        extract_to = rar_path.rsplit('.', 1)[0] + '_extracted'
    
    Path(extract_to).mkdir(parents=True, exist_ok=True)
    
    with rarfile.RarFile(rar_path) as rf:
        rf.extractall(extract_to)
    
    return extract_to


def create_rar(files: List[str], rar_path: str) -> str:
    """Create RAR archive"""
    try:
        import rarfile
    except ImportError:
        raise ConversionError("rarfile required: pip install rarfile")
    
    with rarfile.RarFile(rar_path, 'w') as rf:
        for f in files:
            rf.write(f)
    
    return rar_path


# ============ DATABASE CONVERSIONS ============

def mysql_to_dict(mysql_dump_path: str) -> Dict:
    """Parse MySQL dump file to dictionary"""
    with open(mysql_dump_path, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    
    tables = {}
    current_table = None
    current_data = []
    columns = []
    
    for line in content.split('\n'):
        line = line.strip()
        
        if line.upper().startswith('CREATE TABLE'):
            if current_table and current_data:
                tables[current_table] = {'columns': columns, 'data': current_data}
            match = re.search(r'CREATE TABLE `?(\w+)`?', line, re.IGNORECASE)
            if match:
                current_table = match.group(1)
                current_data = []
                columns = []
        
        elif line.upper().startswith('INSERT INTO'):
            parts = line.split('VALUES')
            if len(parts) > 1:
                values_str = parts[1].rstrip(';')
                if values_str.startswith('(') and values_str.endswith(')'):
                    values_str = values_str[1:-1]
                    values = []
                    in_quote = False
                    current_val = ""
                    for char in values_str:
                        if char == "'" and not in_quote:
                            in_quote = True
                        elif char == "'" and in_quote:
                            in_quote = False
                        elif char == ',' and not in_quote:
                            values.append(current_val.strip().strip("'\""))
                            current_val = ""
                        else:
                            current_val += char
                    values.append(current_val.strip().strip("'\""))
                    
                    if not columns:
                        columns = [f"col_{i}" for i in range(len(values))]
                    current_data.append(dict(zip(columns, values)))
    
    if current_table and current_data:
        tables[current_table] = {'columns': columns, 'data': current_data}
    
    return {'tables': tables}


def postgresql_to_dict(pg_dump_path: str) -> Dict:
    """Parse PostgreSQL dump file to dictionary"""
    with open(pg_dump_path, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    
    tables = {}
    current_table = None
    current_data = []
    columns = []
    
    for line in content.split('\n'):
        line = line.strip()
        
        if line.upper().startswith('CREATE TABLE'):
            if current_table and current_data:
                tables[current_table] = {'columns': columns, 'data': current_data}
            match = re.search(r'CREATE TABLE (\w+)', line, re.IGNORECASE)
            if match:
                current_table = match.group(1)
                current_data = []
                columns = []
        
        elif line.upper().startswith('INSERT INTO'):
            parts = line.split('VALUES')
            if len(parts) > 1:
                values_str = parts[1].rstrip(';')
                if values_str.startswith('(') and values_str.endswith(')'):
                    values_str = values_str[1:-1]
                    values = [v.strip().strip("'\"") for v in values_str.split(',')]
                    if not columns:
                        columns = [f"col_{i}" for i in range(len(values))]
                    current_data.append(dict(zip(columns, values)))
    
    if current_table and current_data:
        tables[current_table] = {'columns': columns, 'data': current_data}
    
    return {'tables': tables}


def dict_to_mysql(data: Dict, table_name: str = "table") -> str:
    """Convert dictionary to MySQL INSERT statements"""
    lines = []
    
    tables = data.get('tables', {table_name: {'columns': list(data.keys()) if not isinstance(data.get('tables'), dict) else [], 'data': [data]}})
    
    for tbl, tbl_data in tables.items():
        cols = tbl_data.get('columns', [])
        rows = tbl_data.get('data', [])
        
        if cols:
            col_defs = ', '.join([f"`{c}` TEXT" for c in cols])
            lines.append(f"CREATE TABLE {tbl} ({col_defs});")
        
        for row in rows:
            if isinstance(row, dict):
                vals = "', '".join(str(v) for v in row.values())
                lines.append(f"INSERT INTO {tbl} VALUES ('{vals}');")
    
    return '\n'.join(lines)


def dict_to_postgresql(data: Dict, table_name: str = "table") -> str:
    """Convert dictionary to PostgreSQL INSERT statements"""
    lines = []
    
    tables = data.get('tables', {table_name: {'columns': list(data.keys()) if not isinstance(data.get('tables'), dict) else [], 'data': [data]}})
    
    for tbl, tbl_data in tables.items():
        cols = tbl_data.get('columns', [])
        rows = tbl_data.get('data', [])
        
        if cols:
            col_defs = ', '.join([f"{c} TEXT" for c in cols])
            lines.append(f"CREATE TABLE {tbl} ({col_defs});")
        
        for row in rows:
            if isinstance(row, dict):
                vals = "', '".join(str(v) for v in row.values())
                lines.append(f"INSERT INTO {tbl} VALUES ('{vals}');")
    
    return '\n'.join(lines)


# ============ WEB/API CONVERSIONS ============

def html_to_text(html_path: str, text_path: str = None) -> str:
    """Convert HTML to plain text"""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ConversionError("beautifulsoup4 required: pip install beautifulsoup4")
    
    with open(html_path, 'r', encoding='utf-8', errors='ignore') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
    
    text = soup.get_text(separator='\n', strip=True)
    
    if text_path is None:
        text_path = html_path.rsplit('.', 1)[0] + '.txt'
    
    with open(text_path, 'w', encoding='utf-8') as f:
        f.write(text)
    
    return text_path


def html_to_markdown(html_path: str, md_path: str = None) -> str:
    """Convert HTML to Markdown"""
    try:
        import html2text
    except ImportError:
        raise ConversionError("html2text required: pip install html2text")
    
    with open(html_path, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    
    h = html2text.HTML2Text()
    h.ignore_links = False
    md = h.handle(content)
    
    if md_path is None:
        md_path = html_path.rsplit('.', 1)[0] + '.md'
    
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(md)
    
    return md_path


def xml_to_dict(xml_path: str) -> Dict:
    """Convert XML to dictionary"""
    try:
        import xmltodict
    except ImportError:
        raise ConversionError("xmltodict required: pip install xmltodict")
    
    with open(xml_path, 'r', encoding='utf-8') as f:
        return xmltodict.parse(f.read())


def dict_to_xml(data: Dict, xml_path: str = None, root: str = "root") -> str:
    """Convert dictionary to XML"""
    try:
        import xmltodict
    except ImportError:
        raise ConversionError("xmltodict required: pip install xmltodict")
    
    xml_str = xmltodict.unparse({root: data}, pretty=True)
    
    if xml_path:
        with open(xml_path, 'w', encoding='utf-8') as f:
            f.write(xml_str)
    
    return xml_str


def url_to_json(url: str) -> Dict:
    """Fetch JSON from URL"""
    try:
        import requests
    except ImportError:
        raise ConversionError("requests required: pip install requests")
    
    response = requests.get(url)
    response.raise_for_status()
    return response.json()


def download_file(url: str, output_path: str = None) -> str:
    """Download file from URL"""
    try:
        import requests
    except ImportError:
        raise ConversionError("requests required: pip install requests")
    
    response = requests.get(url)
    response.raise_for_status()
    
    if output_path is None:
        output_path = url.split('/')[-1]
    
    with open(output_path, 'wb') as f:
        f.write(response.content)
    
    return output_path


# ============ SCIENTIFIC DATA FORMATS ============

def numpy_to_list(array_path: str) -> List:
    """Convert NumPy array file to Python list"""
    try:
        import numpy as np
    except ImportError:
        raise ConversionError("numpy required: pip install numpy")
    
    arr = np.load(array_path)
    return arr.tolist()


def list_to_numpy(data: List, output_path: str = None) -> str:
    """Convert Python list to NumPy array file"""
    try:
        import numpy as np
    except ImportError:
        raise ConversionError("numpy required: pip install numpy")
    
    arr = np.array(data)
    
    if output_path is None:
        output_path = "array.npy"
    
    np.save(output_path, arr)
    return output_path


def hdf5_to_dict(h5_path: str) -> Dict:
    """Read HDF5 file to dictionary"""
    try:
        import h5py
    except ImportError:
        raise ConversionError("h5py required: pip install h5py")
    
    result = {}
    
    def visit_items(name, obj):
        if isinstance(obj, h5py.Dataset):
            result[name] = obj[()]
        elif isinstance(obj, h5py.Group):
            result[name] = dict(obj.attrs)
    
    with h5py.File(h5_path, 'r') as f:
        f.visititems(visit_items)
    
    return result


def matlab_to_dict(mat_path: str) -> Dict:
    """Read MATLAB .mat file to dictionary"""
    try:
        from scipy.io import loadmat
    except ImportError:
        raise ConversionError("scipy required: pip install scipy")
    
    return loadmat(mat_path)


def csv_to_parquet(csv_path: str, parquet_path: str = None) -> str:
    """Convert CSV to Parquet"""
    try:
        import pandas as pd
        import pyarrow as pa
        import pyarrow.parquet as pq
    except ImportError:
        raise ConversionError("pandas and pyarrow required: pip install pandas pyarrow")
    
    df = pd.read_csv(csv_path)
    
    if parquet_path is None:
        parquet_path = csv_path.rsplit('.', 1)[0] + '.parquet'
    
    table = pa.Table.from_pandas(df)
    pq.write_table(table, parquet_path)
    
    return parquet_path


def parquet_to_csv(parquet_path: str, csv_path: str = None) -> str:
    """Convert Parquet to CSV"""
    try:
        import pandas as pd
    except ImportError:
        raise ConversionError("pandas required: pip install pandas")
    
    df = pd.read_parquet(parquet_path)
    
    if csv_path is None:
        csv_path = parquet_path.rsplit('.', 1)[0] + '.csv'
    
    df.to_csv(csv_path, index=False)
    return csv_path


# ============ SERIALIZATION ============

def to_json_bytes(data: Any) -> bytes:
    """Serialize to JSON bytes"""
    return json.dumps(data, default=str).encode('utf-8')


def from_json_bytes(data: bytes) -> Any:
    """Deserialize from JSON bytes"""
    return json.loads(data.decode('utf-8'))


# ============ ENCRYPTION/ENCODING ============

def text_to_qr(text: str, qr_path: str = None) -> str:
    """Generate QR code from text"""
    try:
        import qrcode
    except ImportError:
        raise ConversionError("qrcode required: pip install qrcode")
    
    if qr_path is None:
        qr_path = "qrcode.png"
    
    img = qrcode.make(text)
    img.save(qr_path)
    
    return qr_path


def qr_to_text(qr_path: str) -> str:
    """Read QR code to text"""
    try:
        from PIL import Image
        from pyzbar.pyzbar import decode
    except ImportError:
        raise ConversionError("Pillow and pyzbar required: pip install pillow pyzbar")
    
    img = Image.open(qr_path)
    decoded = decode(img)
    
    if decoded:
        return decoded[0].data.decode('utf-8')
    return ""


def password_to_hash(password: str, algorithm: str = 'bcrypt') -> str:
    """Convert password to hash"""
    if algorithm == 'bcrypt':
        try:
            import bcrypt
            return bcrypt.hashpw(password.encode(), bcrypt.gensalt()).decode()
        except ImportError:
            raise ConversionError("bcrypt required: pip install bcrypt")
    elif algorithm == 'argon2':
        try:
            import argon2
            return argon2.PasswordHasher().hash(password)
        except ImportError:
            raise ConversionError("argon2 required: pip install argon2")
    else:
        return hash_string(password, algorithm)


def verify_password(password: str, hashed: str, algorithm: str = 'bcrypt') -> bool:
    """Verify password against hash"""
    if algorithm == 'bcrypt':
        try:
            import bcrypt
            return bcrypt.checkpw(password.encode(), hashed.encode())
        except ImportError:
            raise ConversionError("bcrypt required: pip install bcrypt")
    else:
        return hash_string(password, algorithm) == hashed


def encode_jwt(payload: Dict, secret: str, algorithm: str = 'HS256') -> str:
    """Encode JWT token"""
    try:
        import jwt
    except ImportError:
        raise ConversionError("PyJWT required: pip install PyJWT")
    
    return jwt.encode(payload, secret, algorithm=algorithm)


def decode_jwt(token: str, secret: str, algorithm: str = 'HS256') -> Dict:
    """Decode JWT token"""
    try:
        import jwt
    except ImportError:
        raise ConversionError("PyJWT required: pip install PyJWT")
    
    return jwt.decode(token, secret, algorithms=[algorithm])


# ============ EBOOK CONVERSIONS ============

def epub_to_text(epub_path: str, txt_path: str = None) -> str:
    """Convert EPUB to plain text"""
    try:
        import ebooklib
        from bs4 import BeautifulSoup
    except ImportError:
        raise ConversionError("ebooklib and beautifulsoup4 required: pip install ebooklib beautifulsoup4")
    
    from ebooklib import epub
    
    book = epub.read_epub(epub_path)
    text_content = []
    
    for item in book.get_items():
        if item.get_type() == 9:
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            text_content.append(soup.get_text(separator='\n', strip=True))
    
    text = '\n\n'.join(text_content)
    
    if txt_path is None:
        txt_path = epub_path.rsplit('.', 1)[0] + '.txt'
    
    with open(txt_path, 'w', encoding='utf-8') as f:
        f.write(text)
    
    return txt_path


def text_to_epub(txt_path: str, epub_path: str = None, title: str = "Book", author: str = "Unknown") -> str:
    """Convert plain text to EPUB"""
    try:
        import ebooklib.epub as epub
    except ImportError:
        raise ConversionError("ebooklib required: pip install ebooklib")
    
    with open(txt_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    book = epub.EpubBook()
    book.set_identifier('id123456')
    book.set_title(title)
    book.set_language('en')
    book.add_author(author)
    
    chapters = content.split('\n\n')
    spine = ['nav']
    
    for i, chapter in enumerate(chapters[:20]):
        if chapter.strip():
            c = epub.EpubHtml(title=f'Chapter {i+1}', file_name=f'chapter_{i+1}.xhtml', lang='en')
            c.content = f'<html><head></head><body><h1>Chapter {i+1}</h1><p>{chapter.replace(chr(10), "<br>")}</p></body></html>'
            book.add_item(c)
            spine.append(c)
    
    book.toc = tuple(book.items)
    book.spine = spine
    
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    
    if epub_path is None:
        epub_path = txt_path.rsplit('.', 1)[0] + '.epub'
    
    epub.write_epub(epub_path, book, {})
    
    return epub_path


# ============ ADVANCED UNIT CONVERSIONS ============

def convert_units_advanced(value: float, from_unit: str, to_unit: str) -> float:
    """Advanced unit conversion using pint library"""
    try:
        from pint import UnitRegistry
    except ImportError:
        raise ConversionError("pint required: pip install pint")
    
    ureg = UnitRegistry()
    quantity = ureg.Quantity(value, from_unit)
    result = quantity.to(to_unit)
    return result.magnitude


def currency_convert(amount: float, from_currency: str, to_currency: str) -> float:
    """Convert currency using forex-python"""
    try:
        from forex_python.converter import CurrencyRates
    except ImportError:
        raise ConversionError("forex-python required: pip install forex-python")
    
    c = CurrencyRates()
    return c.convert(from_currency, to_currency, amount)


def timestamp_to_date(timestamp: Union[int, float], timezone: str = None) -> str:
    """Convert timestamp to date with timezone"""
    from datetime import datetime, timezone as tz
    
    if timezone:
        tz_obj = tz.timezone(timezone)
        dt = datetime.fromtimestamp(timestamp, tz_obj)
    else:
        dt = datetime.fromtimestamp(timestamp)
    
    return dt.strftime('%Y-%m-%d %H:%M:%S %Z')


def date_to_timestamp(date_str: str, format_str: str = '%Y-%m-%d %H:%M:%S') -> float:
    """Convert date string to timestamp"""
    from datetime import datetime
    
    dt = datetime.strptime(date_str, format_str)
    return dt.timestamp()


# ============ IMAGE OCR ============

def image_to_text(image_path: str, language: str = 'eng') -> str:
    """Extract text from image using OCR"""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        raise ConversionError("pytesseract and Pillow required: pip install pytesseract pillow")
    
    img = Image.open(image_path)
    text = pytesseract.image_to_string(img, lang=language)
    return text


# ============ PDF TO IMAGE ============

def pdf_to_images(pdf_path: str, output_dir: str = None, dpi: int = 150) -> List[str]:
    """Convert PDF pages to images"""
    try:
        from pdf2image import convert_from_path
    except ImportError:
        raise ConversionError("pdf2image required: pip install pdf2image")
    
    images = convert_from_path(pdf_path, dpi=dpi)
    
    if output_dir is None:
        output_dir = pdf_path.rsplit('.', 1)[0] + '_pages'
    
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    output_paths = []
    for i, img in enumerate(images):
        output_path = f"{output_dir}/page_{i+1}.png"
        img.save(output_path)
        output_paths.append(output_path)
    
    return output_paths


# ============ ENCRYPTED ZIP ============

def create_encrypted_zip(files: List[str], zip_path: str, password: str) -> str:
    """Create encrypted ZIP archive"""
    try:
        import pyzipper
    except ImportError:
        raise ConversionError("pyzipper required: pip install pyzipper")
    
    with pyzipper.AESZipFile(zip_path, 'w', compression=pyzipper.ZIP_DEFLATED, encryption=pyzipper.WZ_AES) as zf:
        zf.setpassword(password.encode())
        for f in files:
            zf.write(f)
    
    return zip_path


def extract_encrypted_zip(zip_path: str, password: str, extract_to: str = None) -> str:
    """Extract encrypted ZIP archive"""
    try:
        import pyzipper
    except ImportError:
        raise ConversionError("pyzipper required: pip install pyzipper")
    
    if extract_to is None:
        extract_to = zip_path.rsplit('.', 1)[0] + '_extracted'
    
    Path(extract_to).mkdir(parents=True, exist_ok=True)
    
    with pyzipper.AESZipFile(zip_path, 'r') as zf:
        zf.setpassword(password.encode())
        zf.extractall(extract_to)
    
    return extract_to


# ============ MAIN CLI ============

def main():
    import argparse
    
    parser = argparse.ArgumentParser(description='Universal Converter - Convert Anything to Anything')
    parser.add_argument('input', help='Input file, data, or directory')
    parser.add_argument('-o', '--output', help='Output file or directory')
    parser.add_argument('-f', '--from', dest='from_fmt', help='Input format')
    parser.add_argument('-t', '--to', dest='to_fmt', help='Output format')
    parser.add_argument('-l', '--list', action='store_true', help='List supported formats')
    parser.add_argument('-i', '--info', action='store_true', help='Show file info')
    parser.add_argument('--validate', action='store_true', help='Validate file format')
    parser.add_argument('--hash', choices=['md5', 'sha1', 'sha256', 'sha512'], help='File hash')
    parser.add_argument('--resize', action='store_true', help='Resize image')
    parser.add_argument('--width', type=int, help='Target width')
    parser.add_argument('--height', type=int, help='Target height')
    parser.add_argument('--scale', type=float, help='Scale factor (e.g., 0.5)')
    parser.add_argument('--batch', nargs='+', help='Batch convert')
    parser.add_argument('--output-dir', help='Output directory')
    
    args = parser.parse_args()
    
    if args.list:
        formats = get_formats()
        print("Supported Formats:")
        for cat, fmts in formats.items():
            print(f"\n{cat.upper()} ({len(fmts)}):")
            print("  " + ", ".join(fmts))
        return
    
    if args.info:
        print(json.dumps(file_info(args.input), indent=2))
        return
    
    if args.validate:
        print(json.dumps(validate_file(args.input), indent=2))
        return
    
    if args.hash:
        print(f"{args.hash}: {hash_file(args.input, args.hash)}")
        return
    
    if args.resize:
        result = resize_image(args.input, args.output, args.width, args.height, args.scale)
        print(f"Resized: {result}")
        return
    
    if args.batch:
        results = batch_convert(args.batch, args.output_dir, args.to_fmt or 'json')
        for r in results:
            status = "✓" if r['success'] else "✗"
            print(f"{status} {r['file']} → {r.get('output') or r.get('error')}")
        return
    
    try:
        result = convert_file(args.input, args.output, args.from_fmt, args.to_fmt)
        print(f"Converted: {result}")
    except ConversionError as e:
        print(f"Error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
